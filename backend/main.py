from fastapi import FastAPI
from pydantic import BaseModel
import subprocess
import speech_recognition as sr
from gtts import gTTS
import playsound
import pyautogui
from pptx import Presentation
import os

app = FastAPI()

# --- System Prompt ---
system_prompt = (
    "You are an AI teaching assistant for classrooms. "
    "Answer clearly, step by step, using simple examples. "
    "Support the teacher, never reference the internet. "
    "Keep answers short and helpful (max 50 words)."
)

# --- Request Models ---
class AskRequest(BaseModel):
    userPrompt: str
    slide_topic: str | None = None
    slide_num: int | None = None

class VoiceRequest(BaseModel):
    pass  # No device_index, always use default mic


# --- Function to query Ollama ---
def ask_ollama(prompt: str, slide_topic=None, slide_num=None) -> str:
    if slide_topic and slide_num:
        final_prompt = (
            f"{system_prompt}\n\n"
            f"The user asked a question related to slide {slide_num} about {slide_topic}. "
            f"Answer based on that slide’s topic.\n\n{prompt}"
        )
    else:
        final_prompt = system_prompt + "\n\n" + prompt

    result = subprocess.run(
        ["ollama", "run", "qwen2.5:7b"],
        input=final_prompt.encode(),
        capture_output=True
    )
    return result.stdout.decode().strip()


# --- Natural TTS with gTTS ---
def speak(text: str):
    file = "speech.mp3"
    tts = gTTS(text=text, lang="en")
    tts.save(file)
    playsound.playsound(file)
    os.remove(file)


# --- Speech Recognition ---
recognizer = sr.Recognizer()

def listen():
    try:
        with sr.Microphone() as source:  # Use default microphone
            recognizer.adjust_for_ambient_noise(source, duration=1)
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=5)
        return recognizer.recognize_google(audio)
    except Exception as e:
        return None


# --- Slide Control ---
def next_slide():
    pyautogui.press("right")

def previous_slide():
    pyautogui.press("left")

def go_to_slide(n):
    pyautogui.typewrite(str(n))
    pyautogui.press("enter")


# --- Index slides for topic-based navigation ---
def index_slides(file):
    prs = Presentation(file)
    slide_index = {}
    for i, slide in enumerate(prs.slides, start=1):
        text = " ".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        ).lower()
        slide_index[i] = text
    return slide_index

SLIDES = index_slides("lesson.pptx")

def find_slide_by_topic(user_input):
    words = user_input.lower().split()
    for num, text in SLIDES.items():
        for word in words:
            if len(word) > 3 and word in text:
                return num
    return None


# --- API Endpoints ---

@app.post("/text")
def text_endpoint(request: AskRequest):
    """Handles text input, returns JSON only (no TTS)."""
    response = ask_ollama(
        request.userPrompt,
        slide_topic=request.slide_topic,
        slide_num=request.slide_num
    )
    return {"input": request.userPrompt, "response": response}


@app.post("/voice")
def voice_endpoint(request: VoiceRequest):
    """Handles voice input, speaks aloud and returns JSON."""
    spoken_text = listen()
    if not spoken_text:
        return {"response": "No speech detected."}

    response = ask_ollama(spoken_text)
    # speak(response)  # 👈 only /voice triggers TTS
    return {"input": spoken_text, "response": response}
