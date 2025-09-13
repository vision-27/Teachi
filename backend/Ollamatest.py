import subprocess
import speech_recognition as sr
import pyttsx3
import pyautogui
from pptx import Presentation

# --- System Prompt ---
system_prompt = (
    "You are an AI teaching assistant for classrooms. "
    "Answer clearly, step by step, using simple examples. "
    "Support the teacher, never reference the internet. "
    "Keep answers short and helpful (max 50 words)."
)

# --- Function to query Ollama ---
def ask_ollama(prompt, slide_topic=None, slide_num=None):
    if slide_topic and slide_num:
        final_prompt = (
            f"{system_prompt}\n\n"
            f"The user asked a question related to slide {slide_num} about {slide_topic}. "
            f"Answer based on that slide’s topic.\n\n{prompt}"
        )
    else:
        final_prompt = system_prompt + "\n\n" + prompt

    result = subprocess.run(
        ["ollama", "run", "qwen2.5:7b"],
        input=final_prompt.encode(),
        capture_output=True
    )
    return result.stdout.decode().strip()

# --- Setup Speech Recognition ---
recognizer = sr.Recognizer()

def listen():
    try:
        with sr.Microphone() as source:  # Use default microphone
            recognizer.adjust_for_ambient_noise(source, duration=1)
            print("\n🎙 Speak now (5s timeout)...")
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=5)
    except sr.WaitTimeoutError:
        print("⏱️ No speech detected.")
        return None
    except Exception as e:
        print(f"⚠️ Microphone error: {e}")
        return None

    try:
        text = recognizer.recognize_google(audio)
        print("You:", text)
        return text
    except sr.UnknownValueError:
        print("❌ Could not understand")
        return None
    except sr.RequestError:
        print("⚠️ Speech service unavailable")
        return None

# --- Setup TTS (persistent robo voice) ---
engine = pyttsx3.init()
engine.setProperty("rate", 170)  # adjust speaking speed
voices = engine.getProperty("voices")
if voices:
    engine.setProperty("voice", voices[0].id)  # use first available voice

def speak(text):
    print("\nAI:", text)
    engine.say(text)
    engine.runAndWait()

# --- Slide Control ---
def next_slide():
    pyautogui.press("right")

def previous_slide():
    pyautogui.press("left")

def go_to_slide(n):
    pyautogui.typewrite(str(n))
    pyautogui.press("enter")

# --- Index slides for topic-based navigation ---
def index_slides(file):
    prs = Presentation(file)
    slide_index = {}
    for i, slide in enumerate(prs.slides, start=1):
        text = " ".join(
            [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        ).lower()
        slide_index[i] = text
    return slide_index

# Load your PowerPoint here 👇
SLIDES = index_slides("lesson.pptx")

def find_slide_by_topic(user_input):
    words = user_input.lower().split()
    for num, text in SLIDES.items():
        for word in words:
            if len(word) > 3 and word in text:  # ignore short words like "the"
                return num
    return None

# --- Main Loop ---
if __name__ == "__main__":
    print("TutorAI ready 🎓 (Mic: Default). Keep your PowerPoint open in slideshow mode!")

    while True:
        try:
            user_input = listen()
            if not user_input:
                continue

            lower_input = user_input.lower()

            if lower_input in ["exit", "quit", "stop", "bye"]:
                speak("Goodbye!")
                break

            # --- Slide Commands ---
            if "next slide" in lower_input:
                next_slide()
                speak("Moving to next slide.")
                continue
            elif "previous slide" in lower_input:
                previous_slide()
                speak("Going back one slide.")
                continue
            elif lower_input.startswith("go to slide about"):
                topic = lower_input.replace("go to slide about", "").strip()
                slide_num = find_slide_by_topic(topic)
                if slide_num:
                    go_to_slide(slide_num)
                    speak(f"Here is the slide about {topic}.")
                else:
                    speak(f"Sorry, I couldn’t find a slide about {topic}.")
                continue
            elif "slide" in lower_input and any(word.isdigit() for word in lower_input.split()):
                for word in lower_input.split():
                    if word.isdigit():
                        go_to_slide(int(word))
                        speak(f"Going to slide {word}.")
                        break
                continue

            # --- Otherwise: Question for AI ---
            slide_num = find_slide_by_topic(user_input)
            if slide_num:
                # Match found → answer based on that slide
                topic_words = [w for w in user_input.split() if w.lower() in SLIDES[slide_num]]
                slide_topic = topic_words[0] if topic_words else "this topic"
                print(f"\n📝 Related to slide {slide_num} ({slide_topic})\n")
                response = ask_ollama(user_input, slide_topic, slide_num)
                speak(f"Based on slide {slide_num} about {slide_topic}, {response}")
            else:
                # General question
                print("\nThinking...\n")
                response = ask_ollama(user_input)
                speak(response)

        except KeyboardInterrupt:
            print("\nChat ended.")
            break
